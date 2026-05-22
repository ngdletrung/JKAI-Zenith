import os
# 🛡️ [SELF-HEALING-DYNAMIC-LOADER]: Tự động phát hiện và cài đặt phụ thuộc nếu thiếu thưa Master
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    import subprocess
    import sys
    import logging
    logger = logging.getLogger("SOVEREIGN_PRESENTER")
    logger.warning("⚠️ [SELF-HEALING]: Phát hiện thiếu thư viện 'python-pptx'. Tự trị khởi động tiến trình cài đặt tự động...")
    try:
        subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
        logger.info("✅ [SELF-HEALING]: Cài đặt 'python-pptx' thành công thưa Master! Tiến hành nạp lại...")
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
    except Exception as e:
        logger.error(f"❌ [SELF-HEALING-FAILED]: Không thể tự cài đặt phụ thuộc: {e}")
        raise e

# 🛰️ JKAI CORE IMPORTS
from core.utils.engine import engine

class CreativePresenter:
    def __init__(self):
        self.output_dir = os.path.join(os.getenv("WORKSPACE_ROOT", "/app"), "files/Output")
        os.makedirs(self.output_dir, exist_ok=True)
        self.navy_blue = RGBColor(0x1B, 0x3A, 0x57)

    def create_presentation(self, title: str, subtitle: str, slides_content: list, filename: str = "Zenith_Presentation") -> str:
        """Kiến tạo Slide thuyết trình chuyên nghiệp."""
        prs = Presentation()
        output_path = os.path.join(self.output_dir, f"{filename}.pptx")

        # 1. Slide Tiêu đề
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]

        title_shape.text = title.upper()
        subtitle_shape.text = subtitle

        # Định dạng tiêu đề
        for paragraph in title_shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.color.rgb = self.navy_blue
                run.font.bold = True

        # 2. Các Slide nội dung
        for s_data in slides_content:
            slide_layout = prs.slide_layouts[1] # Title and Content
            slide = prs.slides.add_slide(slide_layout)
            
            # Tiêu đề Slide
            title_shape = slide.shapes.title
            title_shape.text = s_data.get("title", "Nội dung")
            
            # Nội dung Slide
            body_shape = slide.placeholders[1]
            tf = body_shape.text_frame
            tf.text = s_data.get("content", "")
            
            # Định dạng nội dung
            for paragraph in tf.paragraphs:
                paragraph.font.size = Pt(18)
                paragraph.space_before = Pt(10)

        prs.save(output_path)
        return output_path

    async def execute_presentation_mission(self, action: str, **kwargs) -> dict:
        """Điều phối các nhiệm vụ thuyết trình."""
        if action == "create_pptx":
            path = self.create_presentation(
                kwargs.get("title"), 
                kwargs.get("subtitle", "Bản thuyết trình từ JKAI Zenith"),
                kwargs.get("slides", []),
                kwargs.get("filename", "Presentation")
            )
            return {"status": "success", "path": path}
        return {"status": "error", "msg": f"Hành động '{action}' chưa được hỗ trợ."}

# 🚀 Singleton
_instance = CreativePresenter()

async def create_slides(**kwargs):
    return await _instance.execute_presentation_mission(action="create_pptx", **kwargs)
